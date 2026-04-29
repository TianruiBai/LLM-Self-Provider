"""Best-effort document text extraction.

All optional dependencies are imported lazily so the provider keeps working
even when none of them are installed; in that case we simply fall back to a
UTF-8 / latin-1 decode of the raw bytes.
"""
from __future__ import annotations

import csv
import io
import json
import os
from typing import Any


# Plain-text-ish file extensions we can decode directly.
_TEXT_EXTS = {
    ".txt", ".md", ".markdown", ".rst", ".log", ".csv", ".tsv", ".json",
    ".jsonl", ".ndjson", ".yml", ".yaml", ".toml", ".ini", ".cfg", ".conf",
    ".html", ".htm", ".xml", ".svg",
    ".py", ".js", ".ts", ".tsx", ".jsx", ".css", ".scss",
    ".c", ".h", ".cc", ".cpp", ".hpp", ".rs", ".go", ".java", ".kt",
    ".rb", ".php", ".sh", ".bash", ".zsh", ".ps1", ".bat", ".sql",
}


def _decode(data: bytes) -> str:
    for enc in ("utf-8", "utf-8-sig", "utf-16", "latin-1"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _from_pdf(data: bytes) -> str:
    try:
        import pypdf  # type: ignore
    except Exception:  # noqa: BLE001
        try:
            from pypdf import PdfReader  # type: ignore
        except Exception:  # noqa: BLE001
            raise RuntimeError("pypdf not installed (pip install pypdf)")
    from pypdf import PdfReader  # type: ignore
    reader = PdfReader(io.BytesIO(data))
    parts: list[str] = []
    for i, page in enumerate(reader.pages):
        try:
            parts.append(f"\n--- page {i + 1} ---\n{page.extract_text() or ''}")
        except Exception:  # noqa: BLE001
            parts.append(f"\n--- page {i + 1} (extract failed) ---\n")
    return "".join(parts).strip()


def _from_docx(data: bytes) -> str:
    try:
        import docx  # type: ignore
    except Exception:  # noqa: BLE001
        raise RuntimeError("python-docx not installed (pip install python-docx)")
    d = docx.Document(io.BytesIO(data))
    parts: list[str] = []
    for p in d.paragraphs:
        if p.text:
            parts.append(p.text)
    for table in d.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            parts.append(" | ".join(cells))
    return "\n".join(parts).strip()


def _from_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except Exception:  # noqa: BLE001
        raise RuntimeError("python-pptx not installed (pip install python-pptx)")
    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides):
        parts.append(f"\n--- slide {i + 1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return "\n".join(parts).strip()


def _from_xlsx(data: bytes) -> str:
    try:
        from openpyxl import load_workbook  # type: ignore
    except Exception:  # noqa: BLE001
        raise RuntimeError("openpyxl not installed (pip install openpyxl)")
    wb = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        parts.append(f"\n--- sheet: {ws.title} ---")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            if any(cells):
                parts.append("\t".join(cells))
    return "\n".join(parts).strip()


def _from_csv(data: bytes) -> str:
    text = _decode(data)
    # Pretty-format as TSV-ish
    out: list[str] = []
    try:
        reader = csv.reader(io.StringIO(text))
        for row in reader:
            out.append("\t".join(row))
        return "\n".join(out)
    except Exception:  # noqa: BLE001
        return text


def _from_rtf(data: bytes) -> str:
    """Strip RTF control codes to plain text without external deps."""
    import re as _re
    text = _decode(data)
    # Drop RTF header / font/colour/stylesheet groups
    text = _re.sub(r"\\\*\\[a-zA-Z]+\s?[^{}]*", " ", text)
    # Convert \uNNNN? escapes to chars (signed 16-bit, may be negative)
    def _u(m: "_re.Match[str]") -> str:
        try:
            cp = int(m.group(1))
            if cp < 0:
                cp += 0x10000
            return chr(cp)
        except Exception:  # noqa: BLE001
            return ""
    text = _re.sub(r"\\u(-?\d+)\??", _u, text)
    # \'XX hex byte
    text = _re.sub(r"\\'([0-9a-fA-F]{2})", lambda m: bytes.fromhex(m.group(1)).decode("latin-1", "replace"), text)
    # paragraph / line breaks
    text = _re.sub(r"\\par[d]?\b", "\n", text)
    text = _re.sub(r"\\line\b", "\n", text)
    text = _re.sub(r"\\tab\b", "\t", text)
    # Drop remaining control words and groups markers
    text = _re.sub(r"\\[a-zA-Z]+-?\d*\s?", "", text)
    text = text.replace("{", "").replace("}", "").replace("\\", "")
    # Collapse excessive whitespace
    text = _re.sub(r"[ \t]+\n", "\n", text)
    text = _re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _from_epub(data: bytes) -> str:
    """Read text from EPUB by unzipping and stripping HTML from content docs."""
    import zipfile
    parts: list[str] = []
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            names = [n for n in zf.namelist() if n.lower().endswith((".xhtml", ".html", ".htm"))]
            names.sort()
            from .tools import html_to_text  # local import; html_to_text is exported
            for n in names:
                try:
                    body = zf.read(n).decode("utf-8", errors="replace")
                except Exception:  # noqa: BLE001
                    continue
                txt = html_to_text(body)
                if txt.strip():
                    parts.append(f"\n--- {n} ---\n{txt.strip()}")
    except Exception as e:  # noqa: BLE001
        raise RuntimeError(f"epub read failed: {e}") from e
    return "\n".join(parts).strip()


def _from_odt(data: bytes) -> str:
    """Pull text out of ODT/ODP/ODS by reading content.xml."""
    import zipfile
    import re as _re
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            try:
                xml = zf.read("content.xml").decode("utf-8", errors="replace")
            except KeyError as e:
                raise RuntimeError("not a valid ODF file (missing content.xml)") from e
    except Exception as e:  # noqa: BLE001
        raise RuntimeError(f"odt read failed: {e}") from e
    # Replace block tags with newlines, then strip remaining tags.
    xml = _re.sub(r"</text:p>", "\n", xml)
    xml = _re.sub(r"</text:h>", "\n", xml)
    xml = _re.sub(r"<text:line-break\s*/?>", "\n", xml)
    xml = _re.sub(r"<text:tab\s*/?>", "\t", xml)
    xml = _re.sub(r"<[^>]+>", "", xml)
    # Decode XML entities
    import html as _html
    xml = _html.unescape(xml)
    xml = _re.sub(r"\n{3,}", "\n\n", xml)
    return xml.strip()


def extract(filename: str, data: bytes) -> dict[str, Any]:
    """Extract plain text from a document blob. Returns {filename, ext, text, format, warning?}."""
    ext = os.path.splitext(filename or "")[1].lower()
    fmt = "text"
    text = ""
    warning: str | None = None
    try:
        if ext == ".pdf":
            fmt = "pdf"
            text = _from_pdf(data)
        elif ext in (".docx",):
            fmt = "docx"
            text = _from_docx(data)
        elif ext in (".pptx",):
            fmt = "pptx"
            text = _from_pptx(data)
        elif ext in (".xlsx", ".xlsm"):
            fmt = "xlsx"
            text = _from_xlsx(data)
        elif ext in (".csv", ".tsv"):
            fmt = "csv"
            text = _from_csv(data)
        elif ext in (".rtf",):
            fmt = "rtf"
            text = _from_rtf(data)
        elif ext in (".epub",):
            fmt = "epub"
            text = _from_epub(data)
        elif ext in (".odt", ".odp", ".ods"):
            fmt = ext.lstrip(".")
            text = _from_odt(data)
        elif ext == ".json" or ext == ".jsonl" or ext == ".ndjson":
            fmt = "json"
            text = _decode(data)
            try:
                # pretty-print json for readability
                parsed = json.loads(text)
                text = json.dumps(parsed, indent=2, ensure_ascii=False)
            except Exception:  # noqa: BLE001
                pass
        elif ext in _TEXT_EXTS or ext == "":
            fmt = "text"
            text = _decode(data)
        else:
            # Last-ditch: try decoding as text anyway.
            fmt = "text"
            text = _decode(data)
            warning = f"Unknown extension '{ext}', decoded as text"
    except Exception as e:  # noqa: BLE001
        warning = f"{type(e).__name__}: {e}"
        text = _decode(data)
    return {
        "filename": filename,
        "ext": ext,
        "format": fmt,
        "text": text,
        "size": len(data),
        **({"warning": warning} if warning else {}),
    }
